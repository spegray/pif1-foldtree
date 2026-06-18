#!/usr/bin/env python3
"""Build the PIF1/RRM3 lab-meeting deck (python-pptx), embedding the verified figures."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
try:
    from PIL import Image
    def aspect(p):
        with Image.open(p) as im: return im.height/im.width
except Exception:
    ASP={"FoldTree_structure_resolves_duplication.png":5.6/13.5,"Q1_saccharomycotina_paralog_retention.png":6.2/11,
         "Q2_pif1_vs_rrm3_identity.png":5.6/13,"tree_schematic.png":6.6/11,"tree_dupclade_real.png":1.0}
    def aspect(p): return ASP.get(os.path.basename(p),0.6)

os.chdir(os.path.expanduser("~/pif1-foldtree"))
FIG="results/figures"
NAVY=RGBColor(0x0A,0x2E,0x3D); WHITE=RGBColor(0xFF,0xFF,0xFF); GOLD=RGBColor(0xE8,0xA3,0x3D)
INK=RGBColor(0x14,0x2A,0x33); MUTE=RGBColor(0x5A,0x6B,0x72); ICE=RGBColor(0xCF,0xE3,0xEC)
BLUE=RGBColor(0x21,0x66,0xAC); RED=RGBColor(0xB2,0x18,0x2B); GREEN=RGBColor(0x1A,0x96,0x50)
HEAD="Cambria"; BODY="Calibri"
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BL=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def slide(bg):
    s=prs.slides.add_slide(BL)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    return s
def box(s,l,t,w,h):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tb,tf
def para(tf,text,size,color,bold=False,font=BODY,align=PP_ALIGN.LEFT,first=False,space=6,bullet=None,italic=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space); p.space_before=Pt(0)
    if bullet:
        rb=p.add_run(); rb.text=bullet+"  "; rb.font.size=Pt(size); rb.font.bold=True
        rb.font.color.rgb=GOLD; rb.font.name=BODY
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return p
def star(s,l,t,size=0.32,color=GOLD):
    sh=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(l),Inches(t),Inches(size),Inches(size))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background(); sh.shadow.inherit=False
    return sh
def fit(s,path,l,t,w,h):
    a=aspect(path); bw,bh=w,w*a
    if bh>h: bh=h; bw=h/a
    s.shapes.add_picture(path,Inches(l+(w-bw)/2),Inches(t+(h-bh)/2),Inches(bw),Inches(bh))
def heading(s,text,sub=None):
    star(s,0.6,0.52)
    _,tf=box(s,1.05,0.42,11.6,1.0)
    para(tf,text,30,INK,bold=True,font=HEAD,first=True,space=2)
    if sub: para(tf,sub,15,MUTE,italic=True)

# ---------- Slide 1: title (dark) ----------
s=slide(NAVY)
star(s,0.85,0.85,0.55)
_,tf=box(s,0.85,2.35,11.6,2.6)
para(tf,"When did PIF1 and RRM3 split?",42,WHITE,bold=True,font=HEAD,first=True,space=10)
para(tf,"Structural phylogenetics places the budding-yeast helicase duplication "
        "at the origin of Saccharomycotina",21,ICE,space=4)
_,tf=box(s,0.9,5.6,11.6,1.2)
para(tf,"957 cellular PIF1-family proteins · 728 fungal species · AlphaFold + foldseek 3Di",15,GOLD,first=True,space=3)
para(tf,"Spencer Gray · 2026-06-18 · github.com/spegray/pif1-foldtree",13,ICE)

# ---------- Slide 2: the question (light) ----------
s=slide(WHITE)
heading(s,"One ancestral gene → two helicases","Why the deep split is hard to see")
_,tf=box(s,0.7,1.85,6.0,5.0)
para(tf,"S. cerevisiae has two PIF1-family 5′→3′ DNA helicases:",17,INK,bold=True,first=True,space=8)
para(tf,"Pif1 — mtDNA maintenance, telomerase inhibition, Okazaki processing",15,INK,bullet="•",space=5)
para(tf,"Rrm3 — replication through protein–DNA barriers (rDNA, tRNA, telomeres)",15,INK,bullet="•",space=12)
para(tf,"Question: where/when did one ancestral fungal PIF1 duplicate into the Pif1 + Rrm3 pair?",16,INK,bold=True,space=12)
para(tf,"Why it's hard: the short, fast-evolving helicase core saturates — amino-acid trees "
        "lose the deep signal and place the split at “Fungi” (kingdom), unresolved.",15,MUTE,italic=True)
fit(s,f"{FIG}/tree_schematic.png",6.9,1.7,6.1,5.4)

# ---------- Slide 3: approach (light) ----------
s=slide(WHITE)
heading(s,"A fungi-wide, structure-aware gene tree")
steps=[("1","957 proteins, 728 species","Cellular PIF1 via InterPro IPR048293 — not the Helitron-contaminated Pfam PF05970."),
       ("2","828 AFDB + 129 ColabFold","AlphaFold structures for all 957; the 129 AFDB gaps folded with ColabFold (AF2). Trimmed to the helicase core."),
       ("3","foldseek 3Di alphabet","Per-residue structural alphabet → an AA+3Di alignment over the same core columns."),
       ("4","reconcile vs species tree","Gene tree → GeneRax duplication–loss reconciliation → the duplication branch (the answer).")]
y=2.0
for n,h,d in steps:
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.8),Inches(y),Inches(0.62),Inches(0.62))
    c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.line.fill.background(); c.shadow.inherit=False
    ctf=c.text_frame; ctf.word_wrap=False; p=ctf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=n; r.font.size=Pt(22); r.font.bold=True; r.font.color.rgb=GOLD; r.font.name=HEAD
    _,tf=box(s,1.7,y-0.06,10.9,1.2)
    para(tf,h,18,INK,bold=True,first=True,space=2)
    para(tf,d,14,MUTE)
    y+=1.28

# ---------- Slide 4: structure resolves (light) ----------
s=slide(WHITE)
heading(s,"The fold carries the signal sequence lost",
        "3Di structural characters — not structural distances — are what resolve the node")
fit(s,f"{FIG}/FoldTree_structure_resolves_duplication.png",0.6,1.95,12.1,4.4)
_,tf=box(s,0.7,6.45,12.0,0.9)
para(tf,"3Di has more informative sites than amino acids (194 vs 180) — yet only the AA+3Di ML "
        "partition reaches Saccharomycotina; pure distance-based FoldTree (fident/alntmscore/lddt) stays deep, like AA-only.",
     13,MUTE,first=True,italic=True)

# ---------- Slide 5: the answer (light) ----------
s=slide(WHITE)
heading(s,"PIF1/RRM3 arose once — in the Saccharomycotina ancestor")
fit(s,f"{FIG}/tree_dupclade_real.png",0.5,1.75,6.6,5.5)
_,tf=box(s,7.4,2.2,5.4,4.8)
para(tf,"The duplication clade: 197 genes / 103 species",17,INK,bold=True,first=True,space=8)
para(tf,"Two mirror-image paralog clades (Pif1 + Rrm3) radiate from a single duplication.",15,INK,bullet="★",space=6)
para(tf,"Both paralogs span ~all Saccharomycotina families → the event predates the budding-yeast radiation.",15,INK,bullet="★",space=6)
para(tf,"Independently confirmed by GeneRax (ML) and species-overlap reconciliation.",15,INK,bullet="★",space=6)
para(tf,"Two-copy mushrooms are a separate, convergent duplication — not yeast Pif1/Rrm3 orthologs.",15,GREEN,bullet="★")

# ---------- Slide 6: reconciliation Q1 (light) ----------
s=slide(WHITE)
heading(s,"Reconciliation pins it to the Saccharomycotina stem")
fit(s,f"{FIG}/Q1_saccharomycotina_paralog_retention.png",0.6,1.8,8.1,5.3)
_,tf=box(s,8.95,2.2,3.9,4.8)
para(tf,"GeneRax (gold-standard ML duplication–loss) maps the duplication to the Saccharomycotina common ancestor.",14,INK,bullet="★",first=True,space=8)
para(tf,"Both paralogs present in 8 of 9 sampled families → stem of the radiation.",14,INK,bullet="★",space=8)
para(tf,"Earliest branch (Lipomycetales) is single-copy; deeper sampling would refine before-vs-after that first split.",14,MUTE,bullet="★",italic=True)

# ---------- Slide 7: Q2 (light) ----------
s=slide(WHITE)
heading(s,"The ancestral fungal PIF1 was “Pif1-like”")
fit(s,f"{FIG}/Q2_pif1_vs_rrm3_identity.png",0.6,1.8,12.1,4.2)
_,tf=box(s,0.7,6.2,12.0,1.1)
para(tf,"A single-copy PIF1 outside the duplication is, by descent, a co-ortholog of both paralogs (0/536 nest inside either). "
        "But by sequence it resembles ScPif1 more (442 vs 93 species; every subphylum) — Pif1 kept more of the ancestral "
        "state; Rrm3 is the derived paralog.",13,MUTE,first=True,italic=True)

# ---------- Slide 8: takeaways (dark) ----------
s=slide(NAVY)
star(s,0.6,0.55,0.4)
_,tf=box(s,1.1,0.5,11.5,0.9); para(tf,"Takeaways",32,WHITE,bold=True,font=HEAD,first=True)
_,tf=box(s,0.9,1.8,11.6,5.2)
for t in [
 "The PIF1/RRM3 duplication maps to the Saccharomycotina ancestor — the stem of the budding-yeast radiation.",
 "Structure (3Di) resolves a node amino acids can't — but as ML characters, not distances.",
 "GeneRax and species-overlap reconciliation agree; the two-copy mushrooms are an independent, convergent duplication.",
 "The ancestral fungal PIF1 was Pif1-like; Rrm3 is the more-derived paralog.",
 "Fully reproducible: github.com/spegray/pif1-foldtree · IQ-TREE 3.1.2, foldseek, GeneRax 2.0.4, ColabFold (AF2).",
]:
    para(tf,t,17,WHITE,bullet="★",first=(t.startswith("The PIF1")),space=14)

out="results/figures/PIF1_RRM3_lab_meeting.pptx"
prs.save(out); print("wrote",out, "slides:",len(prs.slides._sldIdLst))
print("BUILD DECK DONE")
